# -*- coding: utf-8 -*-
"""
A-14-2_noguchi.pptx のスライド7〜11(実験〜まとめ)に、3D姿勢推定(Camera1)の
考察結果を追記するスクリプト。デザインは最小限(白背景・黒文字・罫線のみ)。
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import copy

SRC = "A-14-2_noguchi.pptx"
DST = "A-14-2_noguchi.pptx"

FONT = "Noto Sans JP"
TEXT_COLOR = RGBColor(0x22, 0x22, 0x22)
LEFT = Emu(581025)
WIDTH = Emu(8483399)

PLAIN_TABLE_STYLE = "{5940675A-B579-460E-94D1-54222C63F5DA}"  # No Style, Table Grid


def set_run(run, text, size=16, bold=False, color=TEXT_COLOR, font=FONT):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_bullets(slide, top, height, lines, size=16, width=WIDTH, left=LEFT):
    """lines: list of (text, bold) タプル。先頭が'・'なら箇条書き扱い。"""
    box = slide.shapes.add_textbox(left, Emu(top), width, Emu(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, (text, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        r = p.add_run()
        set_run(r, text, size=size, bold=bold)
    return box


def set_plain_table_style(table):
    tbl = table._tbl
    tblPr = tbl.find(qn('a:tblPr'))
    for child in list(tblPr):
        if child.tag == qn('a:tableStyleId'):
            tblPr.remove(child)
    style_el = tbl.makeelement(qn('a:tableStyleId'), {})
    style_el.text = PLAIN_TABLE_STYLE
    tblPr.append(style_el)
    tblPr.set('firstRow', '1')
    tblPr.set('bandRow', '0')


def add_table(slide, top, rows_data, col_widths, header_size=13, body_size=12, row_h=280000, left=LEFT):
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    width = sum(col_widths)
    height = row_h * n_rows
    gframe = slide.shapes.add_table(n_rows, n_cols, left, Emu(top), Emu(width), Emu(height))
    table = gframe.table
    set_plain_table_style(table)
    for c, w in enumerate(col_widths):
        table.columns[c].width = Emu(w)
    for r in range(n_rows):
        table.rows[r].height = Emu(row_h)
        for c in range(n_cols):
            cell = table.cell(r, c)
            cell.margin_left = Emu(45720)
            cell.margin_right = Emu(45720)
            cell.margin_top = Emu(9144)
            cell.margin_bottom = Emu(9144)
            cell.vertical_anchor = 3  # MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if r == 0 else PP_ALIGN.LEFT
            run = p.add_run()
            is_header = (r == 0)
            set_run(run, str(rows_data[r][c]), size=(header_size if is_header else body_size), bold=is_header)
    return gframe


prs = Presentation(SRC)
slides = list(prs.slides)


# ============================================================
# Slide 7: 恒常性の評価
# ============================================================
s7 = slides[6]
# 既存の方法テキストは残す(工程③ 恒常性の評価 / テスト用データセットで姿勢推定 / 距離を比較)
existing_box = next(sh for sh in s7.shapes if sh.shape_id == 3)
bottom7 = existing_box.top + existing_box.height

add_bullets(s7, bottom7 + 120000, 320000, [
    ("結果: head〜spine5の6区間について、区間長のフレーム間変動係数(CV)を推定・正解で比較(Camera1)", False),
], size=14)

add_table(
    s7, bottom7 + 500000,
    rows_data=[
        ["区間", "推定CV", "正解CV"],
        ["head - spine0", "0.041", "0.042"],
        ["spine0 - spine1", "0.214", "0.000"],
        ["spine1 - spine2", "0.267", "0.003"],
        ["spine2 - spine3", "0.251", "0.000"],
        ["spine3 - spine4", "0.281", "0.000"],
        ["spine4 - spine5", "0.296", "0.000"],
    ],
    col_widths=[3200000, 1600000, 1600000],
    header_size=13, body_size=12, row_h=270000,
)

add_bullets(s7, bottom7 + 500000 + 270000 * 7 + 150000, 900000, [
    ("・正解(剛体骨格)はCVがほぼ0だが、推定はspine系区間で21〜30%変動 → 骨長の恒常性に課題", False),
    ("・head-spine0のみ推定CVが低いのは、モデルのroot関節がx,y=(0,0)に固定される正規化の", False),
    ("　アーティファクトによるもので、精度の高さを意味しない", False),
], size=13)


# ============================================================
# Slide 8: 柔軟性の評価
# ============================================================
s8 = slides[7]

add_bullets(s8, 1533525 + 400000, 500000, [
    ("方法: 肩-肘-手首、股関節-膝-足首の3点から関節角度を算出し、可動範囲・正解との相関を比較(Camera1)", False),
], size=14)

add_table(
    s8, 1533525 + 900000,
    rows_data=[
        ["関節", "推定 可動範囲", "正解 可動範囲", "角度相関"],
        ["left_elbow", "71.4°〜131.0°", "86.7°〜127.4°", "0.90"],
        ["right_elbow", "78.3°〜145.4°", "73.3°〜146.4°", "0.13"],
        ["left_knee", "11.8°〜149.4°", "71.0°〜155.9°", "0.66"],
        ["right_knee", "21.8°〜141.7°", "47.5°〜159.5°", "0.90"],
    ],
    col_widths=[2000000, 2300000, 2300000, 1600000],
    header_size=13, body_size=12, row_h=280000,
)

add_bullets(s8, 1533525 + 900000 + 280000 * 5 + 150000, 900000, [
    ("・left_elbow・right_kneeは動きのタイミングとの相関が高い(0.9前後)", False),
    ("・left_kneeは推定の可動範囲が正解より大きく広がっており、非現実的な曲がり方をしている疑い", False),
], size=13)


# ============================================================
# Slide 9: セルフオクルージョンに対する精度の評価
# ============================================================
s9 = slides[8]
existing_box9 = next(sh for sh in s9.shapes if sh.shape_id == 3)
# 3段落目(旧: 精度的な都合で、評価ができなかった)を方法の説明に差し替え(小さめフォントで1行に収める)
p3 = existing_box9.text_frame.paragraphs[2]
r3 = p3.runs[0]
set_run(r3, "方法: 正解のvisibilityフラグ(遮蔽/可視)で層別し、遮蔽時と可視時の角度誤差を比較(Camera1)", size=14, bold=False)
bottom9 = existing_box9.top + existing_box9.height + 200000

add_table(
    s9, bottom9 + 200000,
    rows_data=[
        ["関節", "遮蔽時MAE", "可視時MAE"],
        ["right_elbow", "24.7° (n=94)", "11.1° (n=7)"],
        ["right_knee", "15.0° (n=24)", "20.9° (n=77)"],
    ],
    col_widths=[2600000, 2600000, 2600000],
    header_size=13, body_size=12, row_h=280000,
)

add_bullets(s9, bottom9 + 200000 + 280000 * 3 + 150000, 1300000, [
    ("・right_elbowは遮蔽時に誤差が明確に悪化 → セルフオクルージョンが精度を下げる典型例", False),
    ("・right_kneeは逆の傾向。遮蔽区間(78〜101フレーム)が動作終盤の低動作量フェーズと重なって", False),
    ("　おり、遮蔽の影響と動作量の影響が交絡している可能性がある", False),
    ("・MotionBertの keypoint_scores は全フレーム常に1.0固定で、自己申告の信頼度は使えなかった", False),
], size=13)


# ============================================================
# Slide 10: 考察
# ============================================================
s10 = slides[9]
existing_box10 = next(sh for sh in s10.shapes if sh.shape_id == 8)
bottom10 = existing_box10.top + existing_box10.height

add_bullets(s10, bottom10 + 150000, 2800000, [
    ("・3D推定モデルは動きのタイミング(方向性)は捉えられているが、剛体骨格としての物理的整合性", False),
    ("　(骨長の恒常性)や関節可動範囲の妥当性には課題が残る", False),
    ("", False),
    ("・セルフオクルージョンは精度に影響するが、動作量など他要因と交絡しやすく、遮蔽の有無だけ", False),
    ("　では単純に評価しきれない", False),
    ("", False),
    ("・本発表の評価はCamera1の単一視点によるものであり、多視点での検証は今後の課題とする", False),
], size=15)


# ============================================================
# Slide 11: まとめ・今後の課題
# ============================================================
s11 = slides[10]
existing_box11 = next(sh for sh in s11.shapes if sh.shape_id == 5)
tf11 = existing_box11.text_frame

# 既存段落: [0]まとめ見出し [1]空 [2]空 [3]空 [4]今後の課題見出し
def fill_para(idx, text, bold=False, size=16):
    p = tf11.paragraphs[idx]
    if p.runs:
        r = p.runs[0]
    else:
        r = p.add_run()
    set_run(r, text, size=size, bold=bold)

fill_para(0, "まとめ", bold=True, size=20)
fill_para(1, "・独自データセットで追加学習したモデルにより3D姿勢推定を実施し、恒常性・柔軟性・セルフ", size=15)
fill_para(2, "　オクルージョン精度の3観点から評価した", size=15)
fill_para(3, "・動作タイミングへの追従は一定確認できたが、骨長の恒常性・関節可動域の妥当性には改善余地", size=15)
fill_para(4, "今後の課題", bold=True, size=20)

# 新しいtextboxで今後の課題の箇条書きを追加(既存boxの直下)
bottom11 = existing_box11.top + existing_box11.height

add_bullets(s11, bottom11 + 150000, 1600000, [
    ("・多視点(全12カメラ)での検証実施", False),
    ("・骨長制約・時間的平滑化の導入によるジッター抑制", False),
    ("・遮蔽に強いモデルにするための学習データ拡充", False),
    ("・動作量とセルフオクルージョンの交絡を排した評価設計", False),
], size=15)

prs.save(DST)
print("done")
